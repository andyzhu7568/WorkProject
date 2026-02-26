import io
import re
from typing import Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from openpyxl import Workbook
from openpyxl.worksheet.worksheet import Worksheet
from openpyxl.worksheet.datavalidation import DataValidation
from openpyxl.styles import Font, PatternFill
from openpyxl.formatting.rule import CellIsRule

YELLOW_FILL = PatternFill(start_color="FFFF00", end_color="FFFF00", fill_type="solid")
ORANGE_FILL = PatternFill(start_color="FFC000", end_color="FFC000", fill_type="solid")
GRAY_FILL = PatternFill(start_color="D9D9D9", end_color="D9D9D9", fill_type="solid")

# Expected/Actual Outcome dropdown: Green, Yellow, Red — font and fill to match source
OUTCOME_GREEN_FONT = Font(color="006100")
OUTCOME_GREEN_FILL = PatternFill(start_color="C6EFCE", end_color="C6EFCE", fill_type="solid")
OUTCOME_YELLOW_FONT = Font(color="9C6500")
OUTCOME_YELLOW_FILL = PatternFill(start_color="FFEB9C", end_color="FFEB9C", fill_type="solid")
OUTCOME_RED_FONT = Font(color="9C0006")
OUTCOME_RED_FILL = PatternFill(start_color="FFC7CE", end_color="FFC7CE", fill_type="solid")

LIST_SHEET_NAME = "_lists"
LOOKUPS_SHEET_NAME = "Lookups"
OUTCOME_NAMED_RANGE = "Outcome"


def _set_outcome_cell(ws: Worksheet, row: int, col: int, value: Optional[str]) -> None:
    """Set cell value and apply Green/Yellow/Red font and fill when value is one of those."""
    cell = ws.cell(row=row, column=col, value=value)
    if value == "Green":
        cell.font = OUTCOME_GREEN_FONT
        cell.fill = OUTCOME_GREEN_FILL
    elif value == "Yellow":
        cell.font = OUTCOME_YELLOW_FONT
        cell.fill = OUTCOME_YELLOW_FILL
    elif value == "Red":
        cell.font = OUTCOME_RED_FONT
        cell.fill = OUTCOME_RED_FILL


def _sanitize_for_excel(text: str) -> str:
    """
    Remove characters that are illegal in Excel (control chars except TAB, LF, CR).
    """
    if not text:
        return text
    allowed = {9, 10, 13}
    return "".join(ch for ch in text if ord(ch) >= 32 or ord(ch) in allowed)


SECTION_TITLE_KEYWORD = "this is the compliance matrix that has been applied to"

# Excel sheet name: max 31 chars, no \ / ? * [ ]
def _sanitize_sheet_name(name: str) -> str:
    s = re.sub(r'[\s\\/?*\[\]:]+', " ", name).strip()[:31]
    return s or "Sheet"


def _sheet_name_from_title(full_title: str, section_index: int) -> str:
    """Derive a sheet name from the slide title after the compliance matrix keyword."""
    lower = (full_title or "").strip().lower()
    if SECTION_TITLE_KEYWORD not in lower:
        return f"Section {section_index}"
    idx = lower.index(SECTION_TITLE_KEYWORD) + len(SECTION_TITLE_KEYWORD)
    suffix = full_title[idx:].strip()
    return _sanitize_sheet_name(suffix) or f"Section {section_index}"


def _init_sheet_headers(ws: Worksheet) -> None:
    """Apply fixed header rows and column widths to a worksheet."""
    ws["A1"] = "Employer"
    ws["A2"] = "Server"
    ws["A3"] = "Test AccountID"
    ws["A4"] = "Config #"

    ws["B5"] = "Question Name: "
    ws["C5"] = "Condition/Response"
    ws["D5"] = "Expected Outcome:"
    ws["E5"] = "Actual Outcome:"
    ws["F5"] = "Comments:"
    ws["G5"] = "Reviewed By:"
    ws["H5"] = "Date"
    ws["I5"] = "Comments"
    for col in "ABCDEFGHI":
        ws[f"{col}5"].fill = YELLOW_FILL

    for col in "ABCDEFGHI":
        cell = ws[f"{col}6"]
        cell.value = None
        cell.fill = ORANGE_FILL

    ws.column_dimensions["A"].width = 20
    ws.column_dimensions["B"].width = 70
    ws.column_dimensions["C"].width = 30
    for col in "DEFGHI":
        ws.column_dimensions[col].width = 20


def _apply_outcome_validation_and_format(ws: Worksheet, first_row: int, last_row: int) -> None:
    """Apply Green/Yellow/Red dropdown and E-column conditional formatting to the given row range."""
    if last_row < first_row:
        return
    # Data validation should only apply when column C has content.
    # Match the source template: list validation uses the named range 'Outcome' (on Lookups sheet)
    # and showDropDown=False so Excel shows the dropdown arrow.
    dv = DataValidation(
        type="list",
        formula1=OUTCOME_NAMED_RANGE,
        allow_blank=True,
        showDropDown=False,
    )

    for r in range(first_row, last_row + 1):
        c_val = ws.cell(row=r, column=3).value
        if c_val is None or str(c_val).strip() == "":
            continue
        dv.add(f"D{r}")
        dv.add(f"E{r}")

    if str(dv.sqref).strip():
        ws.add_data_validation(dv)

    # Conditional formatting for BOTH D and E so user selections show correct font/fill.
    for col_letter in ("D", "E"):
        rng = f"{col_letter}{first_row}:{col_letter}{last_row}"
        ws.conditional_formatting.add(
            rng,
            CellIsRule(
                operator="equal",
                formula=["Green"],
                font=OUTCOME_GREEN_FONT,
                fill=OUTCOME_GREEN_FILL,
            ),
        )
        ws.conditional_formatting.add(
            rng,
            CellIsRule(
                operator="equal",
                formula=["Yellow"],
                font=OUTCOME_YELLOW_FONT,
                fill=OUTCOME_YELLOW_FILL,
            ),
        )
        ws.conditional_formatting.add(
            rng,
            CellIsRule(
                operator="equal",
                formula=["Red"],
                font=OUTCOME_RED_FONT,
                fill=OUTCOME_RED_FILL,
            ),
        )


def _ensure_list_sheet(wb: Workbook) -> None:
    """
    Ensure a hidden worksheet exists with dropdown options:
    A1 = (blank), A2 = Green, A3 = Yellow, A4 = Red.
    """
    if LIST_SHEET_NAME in wb.sheetnames:
        return
    ws = wb.create_sheet(title=LIST_SHEET_NAME)
    ws["A1"] = None
    ws["A2"] = "Green"
    ws["A3"] = "Yellow"
    ws["A4"] = "Red"
    ws.sheet_state = "hidden"


def _ensure_lookups_outcome_named_range(wb: Workbook) -> None:
    """
    Create a hidden 'Lookups' sheet and a named range 'Outcome' that matches the source template
    style, but with the 4 requested states:
      (blank), Green, Yellow, Red
    We place the list in column C to mirror the reference workbook structure.
    """
    if LOOKUPS_SHEET_NAME not in wb.sheetnames:
        ws = wb.create_sheet(title=LOOKUPS_SHEET_NAME)
        ws.sheet_state = "hidden"
    ws = wb[LOOKUPS_SHEET_NAME]

    # Header (optional)
    ws["C2"] = OUTCOME_NAMED_RANGE

    # List values (blank + colors + extra statuses) in C3:C8
    ws["C3"] = None
    ws["C4"] = "Green"
    ws["C5"] = "Yellow"
    ws["C6"] = "Red"
    ws["C7"] = "No Flag"
    ws["C8"] = "N/A"

    # Define or replace named range Outcome → Lookups!$C$3:$C$8
    # openpyxl stores defined names at workbook level.
    if OUTCOME_NAMED_RANGE in wb.defined_names:
        # Remove existing definition(s)
        del wb.defined_names[OUTCOME_NAMED_RANGE]
    wb.create_named_range(OUTCOME_NAMED_RANGE, ws, "$C$3:$C$8")


def _is_grey_empty_row(row) -> bool:
    """Heuristic: row is empty text and has at least one grey-ish filled cell."""
    texts = [cell.text.strip() for cell in row.cells]
    if any(texts):
        return False

    for cell in row.cells:
        try:
            fill = cell.fill
        except AttributeError:
            continue

        if fill is None or fill.type != MSO_FILL_TYPE.SOLID:
            continue

        rgb = getattr(fill.fore_color, "rgb", None)
        if rgb is None:
            continue

        r, g, b = rgb[0], rgb[1], rgb[2]
        # grey-ish and not too dark / not pure white
        if abs(r - g) < 10 and abs(r - b) < 10 and abs(g - b) < 10 and 80 <= r <= 230:
            return True

    return False


def _find_flag_header_row(table) -> Optional[Tuple[int, int]]:
    """
    Return (header_row_idx, flag_col_idx) where 'Flag' appears, or None if not found.
    """
    for r_idx, row in enumerate(table.rows):
        for c_idx, cell in enumerate(row.cells):
            if cell.text.strip().lower() == "flag":
                return r_idx, c_idx
    return None


def _find_condition_columns(table, header_row_idx: int) -> Dict[str, int]:
    """
    Find columns for 'Approved', 'Approved with Restriction', 'Not Approved'
    in the header row and return a mapping: name -> col_idx.
    """
    mapping: Dict[str, int] = {}
    header_row = list(table.rows)[header_row_idx]
    for c_idx, cell in enumerate(header_row.cells):
        text = cell.text.strip().lower()
        if text == "approved":
            mapping["Approved"] = c_idx
        elif text == "approved with restriction":
            mapping["Approved with Restriction"] = c_idx
        elif text == "not approved":
            mapping["Not Approved"] = c_idx
    return mapping


def _append_rows_for_table_row(
    ws: Worksheet,
    current_row: int,
    question_text: str,
    cond_cols: Dict[str, int],
    row,
) -> int:
    """
    Append Excel rows for a single logical question row in the PPT table.
    - For generic questions: one Excel row per non-empty condition cell, with
      B = question text, C = condition description text, D = color
      (Green/Yellow/Red) based on column (Approved / Approved with Restriction / Not Approved).
    - For Yes/No style questions where the PPT uses 'Has Answered Yes' and
      'Has Answered No or Question Unanswered', generate three rows with
      responses: Unanswered, No, Yes (B = question text, C = response).
    """
    text = _sanitize_for_excel(question_text.strip())
    if not text:
        return current_row

    # Collect raw condition descriptions from this row
    cond_texts: Dict[str, str] = {}
    for name, c_idx in cond_cols.items():
        if c_idx < len(row.cells):
            val = row.cells[c_idx].text.strip()
            if val:
                cond_texts[name] = _sanitize_for_excel(val)

    if not cond_texts:
        return current_row

    approved_txt = cond_texts.get("Approved", "").lower()
    not_approved_txt = cond_texts.get("Not Approved", "").lower()

    # Special handling for Yes/No style questions: C = Unanswered / No / Yes; D = Red for Unanswered/No, Green for Yes
    if (
        "has answered yes" in approved_txt
        and "has answered no" in not_approved_txt
        and "question unanswered" in not_approved_txt
    ):
        for label in ["Unanswered", "No", "Yes"]:
            ws.cell(row=current_row, column=1, value=None)
            ws.cell(row=current_row, column=2, value=text)
            ws.cell(row=current_row, column=3, value=label)
            d_value = "Green" if label == "Yes" else "Red"  # Unanswered and No → Red
            _set_outcome_cell(ws, current_row, 4, d_value)
            # E (Actual Outcome) left blank; dropdown added later for whole range
            current_row += 1
        return current_row

    # Generic case: map each non-empty column to a row, using the
    # descriptive text from the PPT cell in column C and a simple
    # color mapping for Expected Outcome in column D.
    color_map = {
        "Approved": "Green",
        "Approved with Restriction": "Yellow",
        "Not Approved": "Red",
    }

    ordered = ["Approved", "Approved with Restriction", "Not Approved"]
    for name in ordered:
        desc = cond_texts.get(name)
        if not desc:
            continue
        ws.cell(row=current_row, column=1, value=None)
        ws.cell(row=current_row, column=2, value=text)
        ws.cell(row=current_row, column=3, value=desc)
        d_value = color_map.get(name)
        _set_outcome_cell(ws, current_row, 4, d_value)
        # E (Actual Outcome) left blank; dropdown and formatting applied later
        current_row += 1

    return current_row


def _append_grey_separator_row(ws: Worksheet, current_row: int) -> int:
    for col_idx in range(1, 10):  # A=1 .. I=9
        cell = ws.cell(row=current_row, column=col_idx, value=None)
        cell.fill = GRAY_FILL
    return current_row + 1


def _append_note_rows(ws: Worksheet, current_row: int, note_text: str) -> int:
    """
    Append a note as its own row in Excel, putting full note text into column B.
    """
    note_text = note_text.strip()
    note_text = _sanitize_for_excel(note_text)
    if not note_text:
        return current_row
    ws.cell(row=current_row, column=2, value=note_text)
    return current_row + 1


def _append_gate_row(ws: Worksheet, current_row: int, gate_text: str) -> int:
    """
    Insert a fixed 'gate' row between question blocks.
    Only column B is filled; other columns are left blank.
    """
    ws.cell(row=current_row, column=1, value=None)
    ws.cell(row=current_row, column=2, value=_sanitize_for_excel(gate_text))
    return current_row + 1


def _get_slide_section_title(slide) -> str:
    """
    Return the text that contains the section keyword, for section detection and sheet naming.
    Checks the title placeholder first, then any shape that contains the keyword
    (so we detect 'This is the compliance matrix that has been applied to' even when
    it is in a text box rather than the title placeholder).
    """
    kw = SECTION_TITLE_KEYWORD
    if slide.shapes.title is not None:
        t = (slide.shapes.title.text or "").strip()
        if kw in t.lower():
            return t
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        try:
            text = (shape.text or "").strip()
        except Exception:
            continue
        if kw in text.lower():
            return text
    return ""


def _process_slide_into_sheet(ws: Worksheet, slide, current_row: int) -> int:
    """Process one slide's tables and notes into the given worksheet. Return next row index."""
    yes_note = "please note: the following factors only apply if you have answered yes"
    no_note = "please note: the following factors only apply if you have answered no"
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.TABLE:
            continue
        table = shape.table
        header_info = _find_flag_header_row(table)
        if not header_info:
            continue
        header_row_idx, flag_col_idx = header_info
        cond_cols = _find_condition_columns(table, header_row_idx)

        for r_idx, row in enumerate(table.rows):
            if r_idx <= header_row_idx:
                continue
            if _is_grey_empty_row(row):
                current_row = _append_grey_separator_row(ws, current_row)
                continue
            question_text = row.cells[flag_col_idx].text.strip()
            if not question_text:
                continue
            lower_q = question_text.lower()
            if yes_note in lower_q:
                current_row = _append_gate_row(
                    ws, current_row, "If anwered Yes to above queestion"
                )
                continue
            if no_note in lower_q:
                current_row = _append_gate_row(
                    ws, current_row, "If anwered No to above queestion"
                )
                continue
            current_row = _append_rows_for_table_row(
                ws, current_row, question_text, cond_cols, row
            )

    if getattr(slide, "has_notes_slide", False) and slide.notes_slide is not None:
        notes_frame = slide.notes_slide.notes_text_frame
        if notes_frame is not None:
            full_notes = notes_frame.text or ""
            for line in full_notes.splitlines():
                lower = line.lower()
                if yes_note in lower:
                    current_row = _append_gate_row(
                        ws, current_row, "If anwered Yes to above queestion"
                    )
                elif no_note in lower:
                    current_row = _append_gate_row(
                        ws, current_row, "If anwered No to above queestion"
                    )

    return current_row


def pptx_to_test_excel(pptx_bytes: bytes) -> bytes:
    """
    Core conversion function:
    - Each time 'This is the compliance matrix that has been applied to' appears in a slide title,
      start a new test project = new Excel tab. Convert content into that tab until the next
      occurrence or end of the presentation.
    """
    prs = Presentation(io.BytesIO(pptx_bytes))
    wb = Workbook()
    # Keep legacy list sheet for backward compatibility, but use Lookups/Outcome to match template.
    _ensure_list_sheet(wb)
    _ensure_lookups_outcome_named_range(wb)

    current_ws: Optional[Worksheet] = None
    current_row = 7
    section_index = 0
    existing_sheet_names: List[str] = [LOOKUPS_SHEET_NAME, LIST_SHEET_NAME]

    for slide in prs.slides:
        title = _get_slide_section_title(slide)
        is_section_start = SECTION_TITLE_KEYWORD in title.lower()

        if is_section_start:
            # Close previous section: apply validation/format to its data range
            if current_ws is not None and current_row >= 7:
                _apply_outcome_validation_and_format(current_ws, 7, current_row - 1)

            section_index += 1
            sheet_name = _sheet_name_from_title(title, section_index)
            # Avoid duplicate sheet names (Excel requires unique tab names)
            base_name = sheet_name
            n = 1
            while sheet_name in existing_sheet_names:
                n += 1
                sheet_name = f"{base_name} ({n})" if base_name else f"Section {section_index}"
            existing_sheet_names.append(sheet_name)

            if current_ws is None:
                # First section: reuse default sheet
                current_ws = wb.active
                current_ws.title = sheet_name
                _init_sheet_headers(current_ws)
            else:
                current_ws = wb.create_sheet(title=sheet_name)
                _init_sheet_headers(current_ws)
            current_row = 7

        if current_ws is not None:
            current_row = _process_slide_into_sheet(current_ws, slide, current_row)

    # If no section title was ever found, treat whole deck as one section (one tab)
    if current_ws is None:
        current_ws = wb.active
        current_ws.title = "Test Sheet"
        _init_sheet_headers(current_ws)
        current_row = 7
        for slide in prs.slides:
            current_row = _process_slide_into_sheet(current_ws, slide, current_row)

    if current_ws is not None and current_row >= 7:
        _apply_outcome_validation_and_format(current_ws, 7, current_row - 1)

    output = io.BytesIO()
    wb.save(output)
    output.seek(0)
    return output.getvalue()

